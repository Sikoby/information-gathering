"""Extract a `DocumentOutline` from an uploaded .pptx or .pdf file.

Pure functions: bytes in, `DocumentOutline` out. No HTTP, no filesystem
side-effects. Used by the `POST /extract` handler in [__main__.py](__main__.py)
and consumed by the generator's propose prompt to drive the document-presentation
mode (one TOPIC per slide, speaker_notes → private_notes).
"""

from __future__ import annotations

import io
from typing import Literal

from .schemas import DocumentOutline, SlideOutline


DocumentKind = Literal["pptx", "pdf"]

_PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)
_PDF_MIME = "application/pdf"

_MAX_SLIDES = 80


def detect_kind(filename: str, content_type: str | None) -> DocumentKind | None:
    """Return 'pptx' or 'pdf' if recognised, else None."""
    name = filename.lower()
    if name.endswith(".pptx") or content_type == _PPTX_MIME:
        return "pptx"
    if name.endswith(".pdf") or content_type == _PDF_MIME:
        return "pdf"
    return None


def parse(filename: str, kind: DocumentKind, data: bytes) -> DocumentOutline:
    if kind == "pptx":
        slides = _parse_pptx(data)
    else:
        slides = _parse_pdf(data)
    if len(slides) > _MAX_SLIDES:
        slides = slides[:_MAX_SLIDES]
    return DocumentOutline(source_name=filename, kind=kind, slides=slides)


def _parse_pptx(data: bytes) -> list[SlideOutline]:
    from pptx import Presentation  # type: ignore[import-not-found]

    prs = Presentation(io.BytesIO(data))
    out: list[SlideOutline] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = _pptx_slide_title(slide)
        content = _pptx_slide_body(slide, skip_title=title)
        notes = _pptx_slide_notes(slide)
        out.append(
            SlideOutline(
                index=idx,
                title=title,
                content=content,
                speaker_notes=notes,
            )
        )
    return out


def _pptx_slide_title(slide) -> str | None:  # noqa: ANN001
    try:
        title_shape = slide.shapes.title
    except (AttributeError, Exception):  # noqa: BLE001 - python-pptx raises bare exceptions
        title_shape = None
    if title_shape is None:
        return None
    text = (title_shape.text or "").strip()
    return text or None


def _pptx_slide_body(slide, skip_title: str | None) -> str:  # noqa: ANN001
    """Collect non-title text shapes, preserving line breaks."""
    lines: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        # Skip the title shape (we already captured it).
        if shape == getattr(slide.shapes, "title", None):
            continue
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if text and text != skip_title:
                lines.append(text)
    return "\n".join(lines).strip()


def _pptx_slide_notes(slide) -> str | None:  # noqa: ANN001
    if not getattr(slide, "has_notes_slide", False):
        return None
    notes_slide = slide.notes_slide
    frame = getattr(notes_slide, "notes_text_frame", None)
    if frame is None:
        return None
    text = (frame.text or "").strip()
    return text or None


def _parse_pdf(data: bytes) -> list[SlideOutline]:
    from pypdf import PdfReader  # type: ignore[import-not-found]

    reader = PdfReader(io.BytesIO(data))
    out: list[SlideOutline] = []
    for idx, page in enumerate(reader.pages, start=1):
        raw = (page.extract_text() or "").strip()
        title, content = _split_first_line(raw)
        out.append(
            SlideOutline(
                index=idx,
                title=title,
                content=content,
                speaker_notes=None,
            )
        )
    return out


def _split_first_line(raw: str) -> tuple[str | None, str]:
    """Treat the first non-empty line of a PDF page as its title."""
    if not raw:
        return None, ""
    lines = [ln.strip() for ln in raw.splitlines() if ln.strip()]
    if not lines:
        return None, ""
    title = lines[0] if len(lines[0]) <= 160 else None
    body = "\n".join(lines[1:] if title else lines)
    return title, body
